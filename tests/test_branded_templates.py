from __future__ import annotations

import io

import pytest

pytest.importorskip("pptx")
pytest.importorskip("docx")

from infrastructure.documents import to_docx, to_pptx  # noqa: E402
from infrastructure.tools import writable_tools  # noqa: E402

_OUTLINE = "## Problem\n- one\n## Proposal\n- two\n- three"


def test_pptx_builds_on_a_template_and_clears_its_slides(tmp_path):
    from pptx import Presentation

    # A "branded" template that already has a stray slide in it.
    tpl = tmp_path / "template.pptx"
    base = Presentation()
    base.slides.add_slide(base.slide_layouts[0])
    base.save(str(tpl))

    data = to_pptx("Deck", _OUTLINE, template=str(tpl))
    deck = Presentation(io.BytesIO(data))

    # Title + Problem + Proposal == 3, and the template's stray slide is gone.
    assert len(deck.slides) == 3


def test_docx_uses_a_template_base(tmp_path):
    from docx import Document

    tpl = tmp_path / "template.docx"
    Document().save(str(tpl))

    data = to_docx("Report", "# Intro\n\nHello.", template=str(tpl))
    assert Document(io.BytesIO(data)) is not None  # opens as a valid docx


def test_missing_template_falls_back_to_plain(tmp_path):
    data = to_pptx("Deck", _OUTLINE, template=str(tmp_path / "nope.pptx"))
    from pptx import Presentation

    assert len(Presentation(io.BytesIO(data)).slides) == 3  # still works


def test_writable_tools_pick_up_a_branding_template(tmp_path):
    from pptx import Presentation

    branding = tmp_path / "branding"
    branding.mkdir()
    Presentation().save(str(branding / "template.pptx"))

    tools = {t.name: t for t in writable_tools(tmp_path, branding=branding)}
    assert tools["write_powerpoint"]._template is not None
    assert tools["write_word"]._template is None  # no doc template provided

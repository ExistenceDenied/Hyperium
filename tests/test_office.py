from __future__ import annotations

import io

from infrastructure.methodologies.json_methodology_repository import (
    JsonMethodologyRepository,
)
from interfaces.office import parse_blocks, to_docx, to_pptx

SAMPLE = """# Report Title

## Overview

This is a paragraph with **bold** words.

- first point
- second point

## Data

| Name | Role |
|------|------|
| Priya | Sponsor |
| Sam | Analyst |
"""

SLIDES = """# Deck

## Slide one
- point A
- point B
Speaker note: keep it short.

## Slide two
- the only point
"""


def test_parser_recognises_the_block_types():
    kinds = [kind for kind, _ in parse_blocks(SAMPLE)]

    assert "h1" in kinds
    assert "h2" in kinds
    assert "p" in kinds
    assert "bullet" in kinds
    assert "table" in kinds


def test_to_docx_produces_a_readable_document():
    from docx import Document

    data = to_docx("Report Title", SAMPLE)
    document = Document(io.BytesIO(data))

    text = "\n".join(p.text for p in document.paragraphs)
    assert "Report Title" in text
    assert "bold" in text  # the ** markers are gone, the word remains
    assert document.tables  # the table came through as a real table


def test_to_pptx_produces_one_slide_per_heading():
    from pptx import Presentation

    data = to_pptx("Deck", SLIDES)
    presentation = Presentation(io.BytesIO(data))

    titles = [
        slide.shapes.title.text
        for slide in presentation.slides
        if slide.shapes.title is not None
    ]
    # A title slide, then one per '##' heading.
    assert "Slide one" in titles
    assert "Slide two" in titles

    # The speaker note landed in the notes, not on the slide body.
    notes = [
        slide.notes_slide.notes_text_frame.text
        for slide in presentation.slides
        if slide.has_notes_slide
    ]
    assert any("keep it short" in note for note in notes)


def test_methodology_declares_a_slide_format():
    repository = JsonMethodologyRepository()
    training = repository.get("training-design")

    formats = {item.key: item.format for item in training.deliverables}

    assert formats["slide-outline"] == "pptx"
    assert formats["facilitator-guide"] == "docx"

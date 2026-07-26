"""
Render a deliverable's Markdown into the file type a client actually receives.

Word and PowerPoint, via python-docx and python-pptx. Those are an optional
dependency — the unit suite and a plain run do not need them — so the imports
are lazy and a clear message points at the extra when they are missing. A small
Markdown block parser sits in front of both, so headings, lists and tables come
through as real Word styles and real slides rather than literal '##' text.
"""

from __future__ import annotations

import io
import re
from email.message import EmailMessage

_HEADING = re.compile(r"^(#{1,6})\s+(.*)")
_BULLET = re.compile(r"^\s*[-*+]\s+")
_NUMBER = re.compile(r"^\s*\d+\.\s+")
_TABLE_SEP = re.compile(r"^\s*\|?[\s:|-]+\|?\s*$")
_BLOCK_START = re.compile(r"^(#{1,6}\s|\s*[-*+]\s|\s*\d+\.\s|\s*\|)")
_BOLD = re.compile(r"(\*\*[^*]+\*\*)")
_NOTE = re.compile(r"(?i)^(speaker note|note)\s*[:\-]\s*")


class OfficeUnavailable(RuntimeError):
    """The optional office dependency is not installed."""


def _cells(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def _plain(text: str) -> str:
    return re.sub(r"[*`]", "", text)


def parse_blocks(markdown: str) -> list[tuple[str, object]]:
    """Turn Markdown into a flat list of (kind, payload) blocks."""
    lines = markdown.splitlines()
    blocks: list[tuple[str, object]] = []
    i, n = 0, len(lines)

    while i < n:
        line = lines[i].rstrip()

        if not line.strip():
            i += 1
            continue

        heading = _HEADING.match(line)
        if heading:
            level = min(len(heading.group(1)), 3)
            blocks.append((f"h{level}", heading.group(2).strip()))
            i += 1
            continue

        if (
            line.lstrip().startswith("|")
            and i + 1 < n
            and _TABLE_SEP.match(lines[i + 1])
            and "-" in lines[i + 1]
        ):
            rows = [_cells(line)]
            i += 2
            while i < n and lines[i].lstrip().startswith("|"):
                rows.append(_cells(lines[i]))
                i += 1
            blocks.append(("table", rows))
            continue

        if _BULLET.match(line):
            items = []
            while i < n and _BULLET.match(lines[i]):
                items.append(_BULLET.sub("", lines[i]).strip())
                i += 1
            blocks.append(("bullet", items))
            continue

        if _NUMBER.match(line):
            items = []
            while i < n and _NUMBER.match(lines[i]):
                items.append(_NUMBER.sub("", lines[i]).strip())
                i += 1
            blocks.append(("number", items))
            continue

        para = [line]
        i += 1
        while i < n and lines[i].strip() and not _BLOCK_START.match(lines[i]):
            para.append(lines[i].rstrip())
            i += 1
        blocks.append(("p", " ".join(para)))

    return blocks


def _add_inline(paragraph, text: str) -> None:
    for part in _BOLD.split(text):
        if part.startswith("**") and part.endswith("**"):
            paragraph.add_run(part[2:-2]).bold = True
        elif part:
            paragraph.add_run(_plain(part))


def to_docx(title: str, markdown: str) -> bytes:
    """Render Markdown into a .docx document, returned as bytes."""
    try:
        from docx import Document
    except ImportError as error:  # pragma: no cover - exercised via the CLI
        raise OfficeUnavailable(
            "python-docx is not installed. Run: pip install -e '.[office]'"
        ) from error

    document = Document()
    document.add_heading(title, level=0)

    for kind, payload in parse_blocks(markdown):
        if kind in ("h1", "h2", "h3"):
            document.add_heading(str(payload), level=int(kind[1]))
        elif kind == "p":
            _add_inline(document.add_paragraph(), str(payload))
        elif kind == "bullet":
            for item in payload:  # type: ignore[union-attr]
                _add_inline(document.add_paragraph(style="List Bullet"), item)
        elif kind == "number":
            for item in payload:  # type: ignore[union-attr]
                _add_inline(document.add_paragraph(style="List Number"), item)
        elif kind == "table":
            rows: list[list[str]] = payload  # type: ignore[assignment]
            if not rows:
                continue
            table = document.add_table(rows=len(rows), cols=len(rows[0]))
            table.style = "Table Grid"
            for r, row in enumerate(rows):
                for c, value in enumerate(row):
                    if c < len(table.rows[r].cells):
                        table.rows[r].cells[c].text = _plain(value)

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _slidify(blocks) -> list[tuple[str | None, list[str], list[str]]]:
    slides: list[tuple[str | None, list[str], list[str]]] = []
    title: str | None = None
    bullets: list[str] = []
    notes: list[str] = []

    def flush() -> None:
        if title or bullets:
            slides.append((title, bullets[:] or ["(no content)"], notes[:]))

    for kind, payload in blocks:
        if kind in ("h1", "h2", "h3"):
            flush()
            title, bullets, notes = str(payload), [], []
        elif kind == "p":
            text = str(payload)
            if _NOTE.match(text):
                notes.append(_NOTE.sub("", text))
            else:
                bullets.append(_plain(text))
        elif kind in ("bullet", "number"):
            bullets.extend(_plain(item) for item in payload)  # type: ignore[union-attr]
        elif kind == "table":
            bullets.extend(" | ".join(row) for row in payload)  # type: ignore[union-attr]

    flush()
    return slides or [(None, ["(no content)"], [])]


_SUBJECT = re.compile(r"(?i)^\s*subject\s*:\s*(.+)$")


def to_eml(title: str, markdown: str) -> bytes:
    """
    Render a deliverable into an .eml draft an email client can open.

    No dependency — this is stdlib email. A recipient is deliberately left off:
    the file is a draft a person opens, addresses and sends. If the content
    opens with a `Subject:` line it is used; otherwise the deliverable name is.
    The body carries both a plain-text and an HTML alternative, so it renders
    well whether the client shows one or the other.
    """
    subject = title
    body = markdown

    first = markdown.lstrip().splitlines()[0] if markdown.strip() else ""
    match = _SUBJECT.match(first)
    if match:
        subject = match.group(1).strip()
        body = markdown.lstrip()[len(first) :].lstrip()

    message = EmailMessage()
    message["Subject"] = subject
    message.set_content(body)

    from interfaces.web.markdown import render

    message.add_alternative(
        f"<html><body>{render(body)}</body></html>", subtype="html"
    )

    return message.as_bytes()


def to_pptx(title: str, markdown: str) -> bytes:
    """Render a slide-outline Markdown into a .pptx deck, returned as bytes."""
    try:
        from pptx import Presentation
    except ImportError as error:  # pragma: no cover - exercised via the CLI
        raise OfficeUnavailable(
            "python-pptx is not installed. Run: pip install -e '.[office]'"
        ) from error

    presentation = Presentation()

    opening = presentation.slides.add_slide(presentation.slide_layouts[0])
    opening.shapes.title.text = title

    for slide_title, bullets, notes in _slidify(parse_blocks(markdown)):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_title or title

        frame = slide.placeholders[1].text_frame
        frame.clear()
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet

        if notes:
            slide.notes_slide.notes_text_frame.text = "\n".join(notes)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()

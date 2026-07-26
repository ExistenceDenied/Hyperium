"""
Render Markdown into the file types a client actually receives: Word and
PowerPoint.

Kept here in infrastructure (not interfaces) so agent tools can produce these
deliverables directly. python-docx and python-pptx are an optional dependency —
the unit suite and a plain run do not need them — so the imports are lazy and a
clear message points at the extra when they are missing. A small Markdown block
parser sits in front of both, so headings, lists and tables come through as real
Word styles and real slides rather than literal '##' text.
"""

from __future__ import annotations

import io
import re
from pathlib import Path

_HEADING = re.compile(r"^(#{1,6})\s+(.*)")
_BULLET = re.compile(r"^\s*[-*+]\s+")
_NUMBER = re.compile(r"^\s*\d+\.\s+")
_TABLE_SEP = re.compile(r"^\s*\|?[\s:|-]+\|?\s*$")
_BLOCK_START = re.compile(r"^(#{1,6}\s|\s*[-*+]\s|\s*\d+\.\s|\s*\|)")
_BOLD = re.compile(r"(\*\*[^*]+\*\*)")
_NOTE = re.compile(r"(?i)^(speaker note|note)\s*[:\-]\s*")


class OfficeUnavailable(RuntimeError):
    """The optional office dependency is not installed."""


def _cells(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def _plain(text: str) -> str:
    return re.sub(r"[*`]", "", text)


def parse_blocks(markdown: str) -> list[tuple[str, object]]:
    """Turn Markdown into a flat list of (kind, payload) blocks."""
    lines = markdown.splitlines()
    blocks: list[tuple[str, object]] = []
    i, n = 0, len(lines)

    while i < n:
        line = lines[i].rstrip()

        if not line.strip():
            i += 1
            continue

        heading = _HEADING.match(line)
        if heading:
            level = min(len(heading.group(1)), 3)
            blocks.append((f"h{level}", heading.group(2).strip()))
            i += 1
            continue

        if (
            line.lstrip().startswith("|")
            and i + 1 < n
            and _TABLE_SEP.match(lines[i + 1])
            and "-" in lines[i + 1]
        ):
            rows = [_cells(line)]
            i += 2
            while i < n and lines[i].lstrip().startswith("|"):
                rows.append(_cells(lines[i]))
                i += 1
            blocks.append(("table", rows))
            continue

        if _BULLET.match(line):
            items = []
            while i < n and _BULLET.match(lines[i]):
                items.append(_BULLET.sub("", lines[i]).strip())
                i += 1
            blocks.append(("bullet", items))
            continue

        if _NUMBER.match(line):
            items = []
            while i < n and _NUMBER.match(lines[i]):
                items.append(_NUMBER.sub("", lines[i]).strip())
                i += 1
            blocks.append(("number", items))
            continue

        para = [line]
        i += 1
        while i < n and lines[i].strip() and not _BLOCK_START.match(lines[i]):
            para.append(lines[i].rstrip())
            i += 1
        blocks.append(("p", " ".join(para)))

    return blocks


def _add_inline(paragraph, text: str) -> None:
    for part in _BOLD.split(text):
        if part.startswith("**") and part.endswith("**"):
            paragraph.add_run(part[2:-2]).bold = True
        elif part:
            paragraph.add_run(_plain(part))


def to_docx(title: str, markdown: str, template=None) -> bytes:
    """
    Render Markdown into a .docx document, returned as bytes.

    A template (.docx/.dotx) is used as the base when given, so the document
    inherits its styles, fonts, header and logo — on-brand rather than blank.
    """
    try:
        from docx import Document
    except ImportError as error:  # pragma: no cover - exercised via the CLI
        raise OfficeUnavailable(
            "python-docx is not installed. Run: pip install -e '.[office]'"
        ) from error

    use = str(template) if template and Path(template).is_file() else None
    document = Document(use) if use else Document()
    document.add_heading(title, level=0)

    for kind, payload in parse_blocks(markdown):
        if kind in ("h1", "h2", "h3"):
            document.add_heading(str(payload), level=int(kind[1]))
        elif kind == "p":
            _add_inline(document.add_paragraph(), str(payload))
        elif kind == "bullet":
            for item in payload:  # type: ignore[union-attr]
                _add_inline(document.add_paragraph(style="List Bullet"), item)
        elif kind == "number":
            for item in payload:  # type: ignore[union-attr]
                _add_inline(document.add_paragraph(style="List Number"), item)
        elif kind == "table":
            rows: list[list[str]] = payload  # type: ignore[assignment]
            if not rows:
                continue
            table = document.add_table(rows=len(rows), cols=len(rows[0]))
            table.style = "Table Grid"
            for r, row in enumerate(rows):
                for c, value in enumerate(row):
                    if c < len(table.rows[r].cells):
                        table.rows[r].cells[c].text = _plain(value)

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _slidify(blocks) -> list[tuple[str | None, list[str], list[str]]]:
    slides: list[tuple[str | None, list[str], list[str]]] = []
    title: str | None = None
    bullets: list[str] = []
    notes: list[str] = []

    def flush() -> None:
        if title or bullets:
            slides.append((title, bullets[:] or ["(no content)"], notes[:]))

    for kind, payload in blocks:
        if kind in ("h1", "h2", "h3"):
            flush()
            title, bullets, notes = str(payload), [], []
        elif kind == "p":
            text = str(payload)
            if _NOTE.match(text):
                notes.append(_NOTE.sub("", text))
            else:
                bullets.append(_plain(text))
        elif kind in ("bullet", "number"):
            bullets.extend(_plain(item) for item in payload)  # type: ignore[union-attr]
        elif kind == "table":
            bullets.extend(" | ".join(row) for row in payload)  # type: ignore[union-attr]

    flush()
    return slides or [(None, ["(no content)"], [])]


def read_text(path) -> str:
    """
    Extract the readable text of a produced deliverable, for review.

    Handles PowerPoint and Word (and spreadsheets) so a critic can judge the
    actual content, not just the agent's summary of it. Anything else is read as
    plain text.
    """
    p = Path(path)
    suffix = p.suffix.lower()
    try:
        if suffix in (".pptx", ".potx"):
            from pptx import Presentation

            parts = []
            for slide in Presentation(str(p)).slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        parts.append(shape.text_frame.text)
            return "\n".join(parts)
        if suffix in (".docx", ".dotx"):
            from docx import Document

            return "\n".join(par.text for par in Document(str(p)).paragraphs)
        if suffix in (".xlsx", ".xlsm"):
            from openpyxl import load_workbook

            rows = []
            for sheet in load_workbook(str(p), data_only=True).worksheets:
                for row in sheet.iter_rows(values_only=True):
                    rows.append(
                        " | ".join("" if c is None else str(c) for c in row)
                    )
            return "\n".join(rows)
        return p.read_text(encoding="utf-8", errors="replace")
    except Exception:
        return ""


def _clear_slides(presentation) -> None:
    """Drop any slides from a template, keeping its theme, master and layouts."""
    slide_ids = presentation.slides._sldIdLst
    for slide_id in list(slide_ids):
        slide_ids.remove(slide_id)


def _set_title(slide, text: str) -> None:
    try:
        if slide.shapes.title is not None:
            slide.shapes.title.text = text
    except (KeyError, AttributeError):
        pass


def _body_frame(slide):
    # The content placeholder is usually index 1, but a branded layout may
    # differ; fall back to the first non-title placeholder.
    try:
        return slide.placeholders[1].text_frame
    except KeyError:
        for placeholder in slide.placeholders:
            if placeholder != slide.shapes.title:
                return placeholder.text_frame
    return None


def to_pptx(title: str, markdown: str, template=None) -> bytes:
    """
    Render a slide-outline Markdown into a .pptx deck, returned as bytes.

    A template (.pptx/.potx) is used as the base when given, so the deck inherits
    its theme, fonts, colours and master layouts — on-brand rather than blank.
    Any slides in the template are cleared; its design is kept.
    """
    try:
        from pptx import Presentation
    except ImportError as error:  # pragma: no cover - exercised via the CLI
        raise OfficeUnavailable(
            "python-pptx is not installed. Run: pip install -e '.[office]'"
        ) from error

    use = str(template) if template and Path(template).is_file() else None
    presentation = Presentation(use) if use else Presentation()
    _clear_slides(presentation)

    layouts = presentation.slide_layouts
    opening = presentation.slides.add_slide(layouts[0])
    _set_title(opening, title)

    content_layout = layouts[1] if len(layouts) > 1 else layouts[0]
    for slide_title, bullets, notes in _slidify(parse_blocks(markdown)):
        slide = presentation.slides.add_slide(content_layout)
        _set_title(slide, slide_title or title)

        frame = _body_frame(slide)
        if frame is not None:
            frame.clear()
            for index, bullet in enumerate(bullets):
                paragraph = (
                    frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                )
                paragraph.text = bullet

        if notes:
            slide.notes_slide.notes_text_frame.text = "\n".join(notes)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
